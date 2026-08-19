# =============================================================================
# app/core/extractor.py —— 内容文本提取 + 元数据
#
# 作用：
#   从常见文档 / 图片 / 压缩包 / 音视频中尽力提取纯文本与元数据，作为 ML
#   特征与概述的输入。所有解析器都是尽力而为：库缺失或解析失败返回 None / {}，
#   绝不抛异常影响主流程。
#
# 结构：
#   extract_text(path) -> str | None          # 按扩展名分发解析器（ML 特征用）
#   extract_metadata(path) -> dict            # 概述用（尺寸/页数/时长/成员…）
#   _read_text / _extract_pdf / _extract_docx / _extract_pptx / _extract_xlsx
#   _extract_archive / _extract_image(OCR) / _extract_audio
#   _ocr_ready() / _ocr_lang() / ocr_hint()  # OCR 能力探测（缓存）
# =============================================================================

"""内容文本与元数据提取：尽力而为，失败返回 None / {}。"""

from __future__ import annotations

import shutil
import subprocess
import tarfile
import zipfile
from pathlib import Path

# 纯文本 / 代码类扩展名：直接按 UTF-8 读取（errors=replace 兜底）
TEXT_EXTS: frozenset[str] = frozenset({
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv",
    ".log", ".json", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".java", ".c", ".h", ".cpp", ".go", ".rs",
    ".yml", ".yaml", ".ini", ".cfg", ".toml", ".sql",
})

# 位图图片（可 OCR）；svg 是矢量图不参与
IMAGE_EXTS: frozenset[str] = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
    ".tiff", ".tif", ".heic", ".jfif",
})

# 压缩包
ARCHIVE_EXTS: frozenset[str] = frozenset({
    ".zip", ".tar", ".tgz", ".gz", ".bz2", ".xz",
})

# 音视频（mutagen 读元数据）
AUDIO_VIDEO_EXTS: frozenset[str] = frozenset({
    ".mp3", ".wav", ".flac", ".aac", ".m4a", ".ogg", ".oga",
    ".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v",
})

# 文本截断上限，控制 TF-IDF 特征规模与训练耗时
MAX_TEXT_CHARS = 100_000


# ---- 主入口 ----

def extract_text(path: str | Path) -> str | None:
    """按扩展名分发解析器提取纯文本；失败或不可解析返回 None。"""
    p = Path(path)
    ext = p.suffix.lower()
    if ext in TEXT_EXTS:
        return _read_text(p)
    if ext == ".pdf":
        return _extract_pdf(p)
    if ext == ".docx":
        return _extract_docx(p)
    if ext == ".pptx":
        return _extract_pptx(p)
    if ext == ".xlsx":
        return _extract_xlsx(p)
    if ext in ARCHIVE_EXTS:
        return _extract_archive(p)
    if ext in IMAGE_EXTS:
        return _extract_image(p)
    if ext in AUDIO_VIDEO_EXTS:
        return _extract_audio(p)
    return None


def extract_metadata(path: str | Path) -> dict:
    """提取概述用的元数据；不可读 / 不支持返回 {}。"""
    p = Path(path)
    ext = p.suffix.lower()
    if ext in IMAGE_EXTS:
        return _image_metadata(p)
    if ext == ".pdf":
        return _pdf_metadata(p)
    if ext == ".pptx":
        return _pptx_metadata(p)
    if ext == ".xlsx":
        return _xlsx_metadata(p)
    if ext in AUDIO_VIDEO_EXTS:
        return _audio_metadata(p)
    if ext in ARCHIVE_EXTS:
        return _archive_metadata(p)
    return {}


# ---- 文本解析 ----

def _read_text(path: Path) -> str | None:
    """读取 UTF-8 文本文件；缺失 / 无权限时返回 None。"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            return _clean(fh.read())
    except OSError:
        return None


def _clean(text: str) -> str:
    """折叠空白并截断长度；空文本返回空字符串。"""
    return " ".join(text.split())[:MAX_TEXT_CHARS]


def _extract_pdf(path: Path) -> str | None:
    """用 PyMuPDF 逐页提取文本；库缺失或解析失败返回 None。"""
    try:
        import pymupdf
    except ImportError:  # 兼容旧版本包名 fitz
        try:
            import fitz as pymupdf
        except ImportError:
            return None
    parts: list[str] = []
    try:
        with pymupdf.open(path) as doc:
            for page in doc:
                parts.append(page.get_text())
    except Exception:
        return None
    return _clean("\n".join(parts)) or None


def _extract_docx(path: Path) -> str | None:
    """用 python-docx 提取段落与表格文本；库缺失或解析失败返回 None。"""
    try:
        from docx import Document
    except ImportError:
        return None
    try:
        doc = Document(str(path))
    except Exception:
        return None
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" ".join(cell.text for cell in row.cells))
    return _clean("\n".join(parts)) or None


def _extract_pptx(path: Path) -> str | None:
    """用 python-pptx 提取全部幻灯片文本；库缺失或解析失败返回 None。"""
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(str(path))
    except Exception:
        return None
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
    return _clean("\n".join(parts)) or None


def _extract_xlsx(path: Path) -> str | None:
    """用 openpyxl 提取全部单元格文本；库缺失或解析失败返回 None。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return None
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return None
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if isinstance(cell, str) and cell.strip():
                        parts.append(cell.strip())
    finally:
        wb.close()
    return _clean("\n".join(parts)) or None


def _extract_archive(path: Path) -> str | None:
    """列出压缩包成员名作为文本特征；损坏或非压缩文件返回 None。"""
    try:
        if path.suffix.lower() == ".zip":
            with zipfile.ZipFile(path) as zf:
                names = zf.namelist()
        else:
            with tarfile.open(path) as tf:
                names = tf.getnames()
    except Exception:
        return None
    return _clean(" ".join(names)) or None


def _extract_audio(path: Path) -> str | None:
    """用 mutagen 提取标题 / 艺术家 / 专辑 / 时长；库缺失或解析失败返回 None。"""
    try:
        from mutagen import File
    except ImportError:
        return None
    try:
        audio = File(path, easy=True)
    except Exception:
        return None
    if audio is None:
        return None
    parts: list[str] = []
    try:
        for key in ("title", "artist", "album"):
            values = audio.get(key)
            if values:
                parts.append(str(values[0]))
        if audio.info is not None and hasattr(audio.info, "length"):
            parts.append(f"时长 {_fmt_duration(audio.info.length)}")
    except Exception:
        pass
    return _clean(" ".join(parts)) or None


# ---- 图片 OCR ----

_tesseract_path: str | None | False = None
_ocr_langs_cache: frozenset[str] | None = None


def _find_tesseract() -> str | None:
    """定位 tesseract 二进制（PATH + 常见安装路径兜底）；未找到缓存 False。"""
    global _tesseract_path
    if _tesseract_path is not None:
        return _tesseract_path or None
    found = shutil.which("tesseract")
    if found is None:
        for base in ("/usr/local/bin", "/opt/homebrew/bin", "/usr/bin"):
            candidate = Path(base) / "tesseract"
            if candidate.exists():
                found = str(candidate)
                break
    _tesseract_path = found or False
    return found or None


def _ocr_langs() -> frozenset[str]:
    """已安装的 tesseract 语言包（缓存，一次探测）。"""
    global _ocr_langs_cache
    if _ocr_langs_cache is not None:
        return _ocr_langs_cache
    binary = _find_tesseract()
    if binary is None:
        _ocr_langs_cache = frozenset()
        return _ocr_langs_cache
    try:
        out = subprocess.run(
            [binary, "--list-langs"], capture_output=True, text=True, timeout=5
        ).stdout
        langs = {
            line.strip() for line in out.splitlines()
            if line.strip() and not line.strip().startswith("List")
        }
    except Exception:
        langs = set()
    _ocr_langs_cache = frozenset(langs)
    return _ocr_langs_cache


def _ocr_lang() -> str:
    """选择 OCR 语言：有 chi_sim 则 中英，否则仅英文。"""
    return "chi_sim+eng" if "chi_sim" in _ocr_langs() else "eng"


def ocr_hint() -> str | None:
    """装了 tesseract 但缺中文语言包时返回一次性提示；否则 None。"""
    if _find_tesseract() and "chi_sim" not in _ocr_langs():
        return (
            "提示：未检测到中文 OCR 语言包（chi_sim），中文图片将用英文识别"
            "（macOS: brew install tesseract-lang）"
        )
    return None


def _extract_image(path: Path) -> str | None:
    """用 tesseract OCR 识别图片文字；无 tesseract / pytesseract 或失败返回 None。"""
    try:
        import pytesseract
    except ImportError:
        return None
    binary = _find_tesseract()
    if binary is None:
        return None
    pytesseract.pytesseract.tesseract_cmd = binary
    try:
        text = pytesseract.image_to_string(str(path), lang=_ocr_lang())
    except Exception:
        return None
    return _clean(text) or None


# ---- 元数据 ----

def _fmt_duration(seconds: float) -> str:
    """秒数 → "M:SS" 或 "H:MM:SS"。"""
    total = int(seconds)
    h, rem = divmod(total, 3600)
    m, s = divmod(rem, 60)
    return f"{h}:{m:02d}:{s:02d}" if h else f"{m}:{s:02d}"


def _image_metadata(path: Path) -> dict:
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS as EXIF_TAGS
    except ImportError:
        return {}
    meta: dict = {}
    try:
        with Image.open(path) as im:
            meta["format"] = im.format or ""
            meta["width"], meta["height"] = im.size
            exif = im.getexif()
            for tag_id, value in exif.items():
                name = EXIF_TAGS.get(tag_id)
                if name == "DateTimeOriginal":
                    meta["date"] = str(value)
                elif name == "Make":
                    meta["make"] = str(value)
                elif name == "Model":
                    meta["model"] = str(value)
    except Exception:
        return {}
    return meta


def _pdf_metadata(path: Path) -> dict:
    try:
        import pymupdf
    except ImportError:
        try:
            import fitz as pymupdf
        except ImportError:
            return {}
    try:
        with pymupdf.open(path) as doc:
            return {"pages": doc.page_count}
    except Exception:
        return {}


def _pptx_metadata(path: Path) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {}
    try:
        return {"slides": len(Presentation(str(path)).slides)}
    except Exception:
        return {}


def _xlsx_metadata(path: Path) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {}
    try:
        wb = load_workbook(str(path), read_only=True)
        try:
            return {"sheets": len(wb.sheetnames)}
        finally:
            wb.close()
    except Exception:
        return {}


def _audio_metadata(path: Path) -> dict:
    try:
        from mutagen import File
    except ImportError:
        return {}
    try:
        audio = File(path, easy=True)
    except Exception:
        return {}
    if audio is None:
        return {}
    meta: dict = {}
    try:
        for key in ("title", "artist", "album"):
            values = audio.get(key)
            if values:
                meta[key] = str(values[0])
        if audio.info is not None and hasattr(audio.info, "length"):
            meta["duration"] = _fmt_duration(audio.info.length)
    except Exception:
        pass
    return meta


def _archive_metadata(path: Path) -> dict:
    try:
        if path.suffix.lower() == ".zip":
            with zipfile.ZipFile(path) as zf:
                names = zf.namelist()
        else:
            with tarfile.open(path) as tf:
                names = tf.getnames()
    except Exception:
        return {}
    return {"member_count": len(names), "members": names[:3]}
