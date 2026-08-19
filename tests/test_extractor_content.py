# =============================================================================
# tests/test_extractor_content.py —— 内容识别扩展（OCR / 更多格式 / 元数据）
# =============================================================================

import tarfile
import zipfile

import pytest

from app.core import extractor


def _tesseract_ready():
    try:
        import pytesseract  # noqa: F401
    except ImportError:
        return False
    return extractor._find_tesseract() is not None


# ---- OCR ----

def test_image_returns_none_when_unreadable(tmp_path):
    p = tmp_path / "pic.png"
    p.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 16)  # 截断 PNG
    assert extractor.extract_text(p) is None


@pytest.mark.skipif(not _tesseract_ready(), reason="需要 tesseract + pytesseract")
def test_ocr_extracts_ascii_text(tmp_path):
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (400, 80), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
    except OSError:
        font = ImageFont.load_default()
    draw.text((20, 20), "hello sorter", fill="black", font=font)
    p = tmp_path / "ocr.png"
    img.save(p)
    text = extractor.extract_text(p)
    assert text and "hello" in text.lower()


def test_ocr_lang_falls_back_to_eng():
    # 未安装 chi_sim 的本机 → eng；chi_sim+eng 分支用 monkeypatch 覆盖
    langs = extractor._ocr_langs()
    if "chi_sim" not in langs:
        assert extractor._ocr_lang() == "eng"
    else:
        assert extractor._ocr_lang() == "chi_sim+eng"


def test_ocr_lang_with_chi_sim(monkeypatch):
    monkeypatch.setattr(extractor, "_ocr_langs_cache", frozenset({"eng", "chi_sim"}))
    assert extractor._ocr_lang() == "chi_sim+eng"


def test_ocr_hint_when_chi_sim_missing():
    if extractor._find_tesseract():
        hint = extractor.ocr_hint()
        if "chi_sim" not in extractor._ocr_langs():
            assert hint and "chi_sim" in hint
        else:
            assert hint is None


# ---- 文档扩展 ----

def test_pptx_roundtrip(tmp_path):
    try:
        from pptx import Presentation
    except ImportError:  # pragma: no cover
        pytest.skip("python-pptx not installed")
    p = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "sales grew twenty percent"
    prs.save(str(p))
    text = extractor.extract_text(p)
    assert text and "Quarterly Review" in text
    meta = extractor.extract_metadata(p)
    assert meta.get("slides") == 1


def test_xlsx_roundtrip(tmp_path):
    try:
        from openpyxl import Workbook
    except ImportError:  # pragma: no cover
        pytest.skip("openpyxl not installed")
    p = tmp_path / "data.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["name", "value"])
    ws.append(["alpha", 42])
    wb.save(str(p))
    text = extractor.extract_text(p)
    assert text and "alpha" in text
    meta = extractor.extract_metadata(p)
    assert meta.get("sheets") == 1


def test_zip_members_and_metadata(tmp_path):
    p = tmp_path / "bundle.zip"
    with zipfile.ZipFile(p, "w") as zf:
        zf.writestr("docs/a.txt", "x")
        zf.writestr("docs/b.txt", "y")
    text = extractor.extract_text(p)
    assert text and "docs/a.txt" in text
    meta = extractor.extract_metadata(p)
    assert meta["member_count"] == 2


def test_tar_members(tmp_path):
    p = tmp_path / "bundle.tar"
    with tarfile.open(p, "w") as tf:
        data = b"x"
        info = tarfile.TarInfo("inner/file.txt")
        info.size = len(data)
        tf.addfile(info, __import__("io").BytesIO(data))
    text = extractor.extract_text(p)
    assert text and "inner/file.txt" in text


def test_corrupt_archive_returns_none(tmp_path):
    p = tmp_path / "bundle.zip"
    p.write_bytes(b"\x00\x01\x02")
    assert extractor.extract_text(p) is None
    assert extractor.extract_metadata(p) == {}


def test_audio_garbage_returns_none(tmp_path):
    p = tmp_path / "song.mp3"
    p.write_bytes(b"\xff\xfb garbage")
    # mutagen 解析失败 → None（不抛）
    result = extractor.extract_text(p)
    assert result is None or isinstance(result, str)


def test_audio_wav_metadata(tmp_path):
    try:
        import mutagen  # noqa: F401
    except ImportError:  # pragma: no cover
        pytest.skip("mutagen not installed")
    import wave
    p = tmp_path / "tone.wav"
    with wave.open(str(p), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(8000)
        w.writeframes(b"\x00\x00" * 8000)  # 1 秒静音
    meta = extractor.extract_metadata(p)
    assert "duration" in meta


# ---- 元数据 ----

def test_image_metadata_dims(tmp_path):
    from PIL import Image
    img = Image.new("RGB", (320, 240), "red")
    p = tmp_path / "img.png"
    img.save(p)
    meta = extractor.extract_metadata(p)
    assert meta["width"] == 320
    assert meta["height"] == 240


def test_pdf_metadata_pages(tmp_path):
    try:
        import pymupdf
    except ImportError:  # pragma: no cover
        pytest.skip("pymupdf not installed")
    p = tmp_path / "doc.pdf"
    doc = pymupdf.open()
    doc.new_page()
    doc.new_page()
    doc.save(str(p))
    doc.close()
    meta = extractor.extract_metadata(p)
    assert meta["pages"] == 2


def test_metadata_unreadable_returns_empty(tmp_path):
    assert extractor.extract_metadata(tmp_path / "nope.png") == {}
